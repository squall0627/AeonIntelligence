import os
import math

from enum import Enum


from concurrent.futures import ThreadPoolExecutor
from io import StringIO
from pathlib import Path

import pptx
from pptx.enum.text import MSO_AUTO_SIZE

from core.ai_core.translation.file_translator.file_translator_base import (
    FileTranslatorBase,
)
from core.ai_core.translation.file_translator.file_translator_type import (
    FileTranslatorType,
)
from core.ai_core.translation.language import Language
from core.utils.log_handler import rotating_file_logger

from core.utils.markitdown import PptxConverter

logger = rotating_file_logger("ai_core")

_FONT_NAME = {
    Language.ENGLISH: "Arial",
    Language.JAPANESE: "Meiryo UI",
    Language.CHINESE: "Microsoft YaHei",
}

_FALLBACK_FONT = "Arial"  # A universal fallback font


class TranslationMode(str, Enum):
    TRANSLATE = "1"
    EXTRACT = "2"
    REPLACE = "3"


# Function for parallel translation
def translate_texts(texts, translator):
    translated_texts = [translator.translate(text) for text in texts]
    return translated_texts


class PPTXTranslator(FileTranslatorBase):
    def __init__(self):
        super().__init__(FileTranslatorType.PPTX)

    async def translate_impl(self, output_dir: Path | str) -> Path | str:
        logger.info(
            f"Translating {self.input_file_path} to {output_dir} by PPTXTranslator"
        )

        run_parallely = self.kwargs.get("run_parallely", True)
        target_slide_index = self.kwargs.get("target_slide_index", None)

        text_translator = self.text_translator
        ppt = pptx.Presentation(self.input_file_path)

        logger.info(
            f"translating slides {'parallely' if run_parallely else 'sequentially'}"
        )
        slide_index = 0
        if run_parallely:
            # run parallely
            target_slides = []
            for slide in ppt.slides:
                if target_slide_index is None or slide_index in target_slide_index:
                    target_slides.append(slide)

                slide_index += 1

            # Extract all slide texts
            slides_texts = [
                self._translate_slide(slide, TranslationMode.EXTRACT)
                for slide in target_slides
            ]

            # Translate texts in parallel
            with ThreadPoolExecutor(max_workers=8) as executor:
                futures = {}
                for idx, texts in enumerate(slides_texts):
                    futures[idx] = {
                        "task": executor.submit(
                            translate_texts, texts, text_translator
                        ),
                        "status": "pending",
                    }

                try:
                    # Gather the results
                    while sum(
                        1 for _, future in futures.items() if future["status"] == "done"
                    ) < len(futures.keys()):
                        for idx, future_task in futures.items():
                            if (
                                future_task["task"].done()
                                and future_task["status"] == "pending"
                            ):
                                translated_texts = future_task["task"].result()
                                # Replace the translated content back in slides
                                self._translate_slide(
                                    target_slides[idx],
                                    TranslationMode.REPLACE,
                                    translated_texts,
                                )
                                future_task["status"] = "done"

                                # completion rate compute
                                self.completion_rate = math.floor(
                                    sum(
                                        1
                                        for _, future in futures.items()
                                        if future["status"] == "done"
                                    )
                                    / len(futures.keys())
                                    * 100
                                )
                except Exception as e:
                    logger.error(e)
                    print(e)
        else:
            # run sequentially
            translated_slide_count = 0
            for slide in ppt.slides:
                if target_slide_index is None or slide_index in target_slide_index:
                    self._translate_slide(slide, TranslationMode.TRANSLATE)
                    translated_slide_count += 1
                    # completion rate compute
                    self.completion_rate = math.floor(
                        translated_slide_count
                        / len(
                            target_slide_index
                            if target_slide_index is not None
                            else ppt.slides
                        )
                        * 100
                    )
                slide_index += 1

        logger.debug(">>> Translating input file name")
        # translate the input file name
        input_file_name = Path(self.input_file_path).name
        output_file_name = text_translator.translate(input_file_name)
        output_path = os.path.join(output_dir, output_file_name)

        # save translated file
        ppt.save(output_path)

        logger.info(f"Translated {self.input_file_path} save to {output_path}")

        return output_path

    def _translate_slide(
        self, slide, mode: TranslationMode, translated_texts: list | None = None
    ) -> list:
        text_translator = self.text_translator
        converter = PptxConverter()
        is_translate_picture = self.kwargs.get("is_translate_picture", False)
        extract_texts = []
        text_idx = 0
        for shape in slide.shapes:
            if is_translate_picture:
                logger.debug(">>> Translating Pictures")
                # Pictures
                if converter._is_picture(shape):
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    if mode == TranslationMode.TRANSLATE:
                        alt_text_translated = text_translator.translate(alt_text)
                        shape.image.alt_text = alt_text_translated
                    elif mode == TranslationMode.EXTRACT:
                        extract_texts.append(alt_text)
                    elif mode == TranslationMode.REPLACE:
                        shape.image.alt_text = translated_texts[text_idx]
                        text_idx += 1

            logger.debug(">>> Translating Tables")
            # Tables
            if converter._is_table(shape):
                for row in shape.table.rows:
                    for cell in row.cells:
                        sub_extract_texts, text_idx = self._translate_text_with_style(
                            cell.text_frame,
                            text_translator,
                            mode,
                            translated_texts,
                            text_idx,
                        )
                        extract_texts.extend(sub_extract_texts)

                # adjust font size to adapt to the size of shape
                self._adjust_font_size_to_fit_shape(shape)

            logger.debug(">>> Translating Charts")
            # Charts
            if shape.has_chart:
                chart = shape.chart
                # Translate the chart title
                if chart.has_title and chart.chart_title.has_text_frame:
                    sub_extract_texts, text_idx = self._translate_text_with_style(
                        chart.chart_title.text_frame,
                        text_translator,
                        mode,
                        translated_texts,
                        text_idx,
                    )
                    extract_texts.extend(sub_extract_texts)

            logger.debug(">>> Translating Text Areas")
            # Text areas
            if shape.has_text_frame:
                text_frame = shape.text_frame
                sub_extract_texts, text_idx = self._translate_text_with_style(
                    text_frame, text_translator, mode, translated_texts, text_idx
                )
                extract_texts.extend(sub_extract_texts)
                # adjust font size to adapt to the size of shape
                self._adjust_font_size_to_fit_shape(shape)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                sub_extract_texts, text_idx = self._translate_text_with_style(
                    notes_frame, text_translator, mode, translated_texts, text_idx
                )
                extract_texts.extend(sub_extract_texts)

        return extract_texts

    def _translate_text_with_style(
        self,
        text_frame,
        text_translator,
        mode: TranslationMode,
        translated_texts: list,
        text_idx: int,
    ) -> (list, int):
        """
        Translates the text inside a PowerPoint shape while preserving styling.

        :param text_frame: The text frame of the shape containing text.
        :param text_translator: The text translator object (e.g., translation API or function).
        """
        if text_frame is None:
            return

        # Retrieve the current vertical alignment of the cell's text frame
        vertical_alignment = (
            text_frame.vertical_anchor if text_frame is not None else None
        )

        extract_texts = []
        # Iterate over each paragraph in the text frame
        for paragraph in text_frame.paragraphs:
            alignment = paragraph.alignment
            # Create a list to hold translated runs to reconstruct the paragraph later
            translated_runs = []

            # Iterate over each run in the paragraph
            original_text = StringIO()
            for run in paragraph.runs:
                # original_text = run.text  # Get text in the run
                # translated_text = text_translator.translate(
                #     original_text
                # )  # Translate the text
                # translated_runs.append(
                #     (translated_text, run.font)
                # )  # Save translated text and font
                original_text.write(run.text)

            if original_text.tell() == 0:
                continue

            if mode == TranslationMode.TRANSLATE or mode == TranslationMode.REPLACE:
                translated_text = ""
                if mode == TranslationMode.TRANSLATE:
                    translated_text = text_translator.translate(
                        original_text.getvalue()
                    )
                elif mode == TranslationMode.REPLACE:
                    translated_text = translated_texts[text_idx]
                    text_idx += 1
                translated_runs.append((translated_text, paragraph.runs[0].font))

                # Clear the paragraph and replace it with translated runs
                paragraph.clear()  # Remove existing text in the paragraph
                for translated_text, original_font in translated_runs:
                    new_run = paragraph.add_run()  # Add new run
                    new_run.text = translated_text  # Set the translated text

                    # Manually copy font properties
                    if original_font is not None:  # Ensure the original font exists
                        if original_font.name:
                            new_run.font.name = _FONT_NAME[self.target_language]
                        if original_font.size:
                            new_run.font.size = original_font.size
                        new_run.font.bold = original_font.bold
                        new_run.font.italic = original_font.italic
                        new_run.font.underline = original_font.underline
                        # Safely copy the color if it exists and has the rgb property
                        if original_font.color and hasattr(original_font.color, "rgb"):
                            new_run.font.color.rgb = original_font.color.rgb
                        else:
                            # Fallback to black color if original color is undefined
                            from pptx.dml.color import RGBColor

                            new_run.font.color.rgb = RGBColor(0, 0, 0)

                # If the table cell has specific alignments (horizontal or vertical), ensure they are preserved since clearing and reconstructing text may default some alignments.
                paragraph.alignment = alignment
            elif mode == TranslationMode.EXTRACT:
                extract_texts.append(original_text.getvalue())

        # Restore the vertical alignment after clearing/replacing text
        if mode == TranslationMode.TRANSLATE or mode == TranslationMode.REPLACE:
            if vertical_alignment:
                text_frame.vertical_anchor = vertical_alignment

        return extract_texts, text_idx

    def _adjust_font_size_to_fit_shape(self, shape):
        """
        Adjusts the font size of text in the shape to fit within the shape's dimensions.

        :param shape: A shape object with a text frame (pptx.shapes.base.Shape).
        """
        if not shape.has_text_frame:
            return  # Skip shapes without a text frame

        text_frame = shape.text_frame
        if not text_frame or not text_frame.text:  # Skip if no text exists
            return

        # Set auto-size and word wrap to ensure proper fitting
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True

        # # Set font family safely with fallback
        # font_family = _FONT_NAME.get(self.target_language, _FALLBACK_FONT)
        #
        # try:
        #     # Attempt to fit text using the target font
        #     text_frame.fit_text(font_family=font_family)
        # except KeyError as e:
        #     logger.warning(
        #         f"Font '{font_family}' not found. Falling back to default font '{_FALLBACK_FONT}'."
        #     )
        #     try:
        #         text_frame.fit_text(font_family=_FALLBACK_FONT)
        #     except Exception:
        #         logger.warning(
        #             f"text_frame.fit_text() failed with fallback font '{_FALLBACK_FONT}'. Skipping."
        #         )

        # # Get the maximum dimensions of the shape
        # max_width = shape.width
        # max_height = shape.height
        #
        # fits = False
        # min_font_size = 100
        # while True:
        #     # Measure the text's bounding box after applying the font size
        #     # Note: python-pptx itself does not calculate text dimensions; this is theoretical sizing.
        #     text_width = text_frame._bodyPr.autofit()
        #     text_height = text_frame._bodyPr.txXform.ext.y
        #
        #     # Check if the text fits within the shape's dimensions
        #     if text_width <= max_width and text_height <= max_height:
        #         fits = True
        #         break
        #
        #     if min_font_size <= 6:
        #         break
        #
        #     # Try setting the font size
        #     for paragraph in text_frame.paragraphs:
        #         for run in paragraph.runs:
        #             run.font.size = Pt(run.font.size - 1)
        #             if run.font.size < min_font_size:
        #                 min_font_size = run.font.size
        #
        # if not fits:
        #     logger.warn(
        #         f"Warning: Text does not fit into the shape within the minimum font size of {min_font_size}pt"
        #     )
